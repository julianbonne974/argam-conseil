#!/usr/bin/env python3
"""
Script de création du PowerPoint COGOHR - Argam Conseils
Présentation de 7 slides pour Clémence Gavache
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Couleurs Argam Conseils
GOLDEN_BROWN = RGBColor(0xB4, 0x92, 0x5E)  # #B4925E - Primary
AUBERGINE = RGBColor(0x52, 0x4C, 0x5D)  # #524C5D - Secondary
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
RED = RGBColor(0xDC, 0x26, 0x26)  # Pour les alertes
GREEN = RGBColor(0x10, 0xB9, 0x81)  # Pour le positif
DARK_GRAY = RGBColor(0x64, 0x74, 0x8B)


def add_title_slide(prs, title, subtitle):
    """Ajoute un slide de titre"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Fond doré subtil en haut
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLDEN_BROWN
    shape.line.fill.background()

    # Titre principal
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = AUBERGINE
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2), Inches(12.33), Inches(1)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = GOLDEN_BROWN
    p.alignment = PP_ALIGN.CENTER

    # Présentatrice
    presenter_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.8), Inches(12.33), Inches(0.5)
    )
    tf = presenter_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Présenté par Clémence Gavache"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Badges partenaires
    badge_box = slide.shapes.add_textbox(
        Inches(4), Inches(6.5), Inches(5.33), Inches(0.4)
    )
    tf = badge_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Partenaire Officiel COGOHR"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLDEN_BROWN
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_header(slide, title, slide_num):
    """Ajoute l'en-tête standard à un slide"""
    # Bande dorée en haut
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLDEN_BROWN
    shape.line.fill.background()

    # Titre du slide
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(11), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = AUBERGINE

    # Numéro de slide
    num_box = slide.shapes.add_textbox(Inches(12), Inches(0.3), Inches(1), Inches(0.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_num}/7"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_footer(slide):
    """Ajoute le pied de page standard"""
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.2), Inches(12.33), Inches(0.3)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Argam Conseils x COGOHR | Présentation exclusive adhérents"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER


def create_slide_2_problem(prs):
    """Slide 2: Le Problème - 53% ne comptent pas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "53% de vos revenus ne comptent pas pour la retraite", 2)

    # Icône alerte (rectangle rouge)
    alert = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(1.8), Inches(0.5)
    )
    alert.fill.solid()
    alert.fill.fore_color.rgb = RED
    alert.line.fill.background()
    alert_tf = alert.text_frame
    alert_tf.paragraphs[0].text = "⚠ ALERTE"
    alert_tf.paragraphs[0].font.color.rgb = WHITE
    alert_tf.paragraphs[0].font.bold = True
    alert_tf.paragraphs[0].font.size = Pt(14)
    alert_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Graphique simplifié - Deux barres
    # Barre Prime 53%
    bar1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.2), Inches(5), Inches(1.2)
    )
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = RED
    bar1.line.fill.background()
    bar1_tf = bar1.text_frame
    bar1_tf.paragraphs[0].text = "Prime de vie chère : 53%"
    bar1_tf.paragraphs[0].font.color.rgb = WHITE
    bar1_tf.paragraphs[0].font.bold = True
    bar1_tf.paragraphs[0].font.size = Pt(20)
    bar1_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    bar1_tf.word_wrap = True

    # Label "NON COMPTÉE"
    label1 = slide.shapes.add_textbox(
        Inches(6.2), Inches(2.5), Inches(2.5), Inches(0.6)
    )
    tf = label1.text_frame
    tf.paragraphs[0].text = "❌ NON COMPTÉE"
    tf.paragraphs[0].font.color.rgb = RED
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(16)

    # Barre Traitement 47%
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.6), Inches(4.4), Inches(1.2)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = GREEN
    bar2.line.fill.background()
    bar2_tf = bar2.text_frame
    bar2_tf.paragraphs[0].text = "Traitement indiciaire : 47%"
    bar2_tf.paragraphs[0].font.color.rgb = WHITE
    bar2_tf.paragraphs[0].font.bold = True
    bar2_tf.paragraphs[0].font.size = Pt(20)
    bar2_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    bar2_tf.word_wrap = True

    # Label "COMPTÉE"
    label2 = slide.shapes.add_textbox(
        Inches(5.6), Inches(3.9), Inches(2.5), Inches(0.6)
    )
    tf = label2.text_frame
    tf.paragraphs[0].text = "✓ COMPTÉE"
    tf.paragraphs[0].font.color.rgb = GREEN
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(16)

    # Tableau des chiffres
    table_x = Inches(8.5)
    table_y = Inches(2)

    # Box tableau
    table_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.8), Inches(4.5), Inches(3.2)
    )
    table_bg.fill.solid()
    table_bg.fill.fore_color.rgb = LIGHT_GRAY
    table_bg.line.color.rgb = GOLDEN_BROWN

    # Titre tableau
    tb_title = slide.shapes.add_textbox(
        Inches(8.5), Inches(2), Inches(4.1), Inches(0.5)
    )
    tf = tb_title.text_frame
    tf.paragraphs[0].text = "Exemple concret"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = AUBERGINE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Lignes du tableau
    lines = [
        ("Revenu mensuel", "3 000 €", BLACK),
        ("Prime vie chère (53%)", "1 590 €", RED),
        ("Traitement indiciaire (47%)", "1 410 €", GREEN),
    ]

    y_pos = 2.6
    for label, value, color in lines:
        # Label
        lb = slide.shapes.add_textbox(
            Inches(8.5), Inches(y_pos), Inches(2.8), Inches(0.4)
        )
        tf = lb.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = DARK_GRAY

        # Valeur
        vb = slide.shapes.add_textbox(
            Inches(11.3), Inches(y_pos), Inches(1.3), Inches(0.4)
        )
        tf = vb.text_frame
        tf.paragraphs[0].text = value
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        y_pos += 0.5

    # Message clé en bas
    msg_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(5.3),
        Inches(12.33),
        Inches(1.2),
    )
    msg_box.fill.solid()
    msg_box.fill.fore_color.rgb = AUBERGINE
    msg_box.line.fill.background()

    msg_tf = msg_box.text_frame
    msg_tf.paragraphs[
        0
    ].text = "Seul le traitement indiciaire est pris en compte pour calculer votre pension de retraite."
    msg_tf.paragraphs[0].font.color.rgb = WHITE
    msg_tf.paragraphs[0].font.size = Pt(22)
    msg_tf.paragraphs[0].font.bold = True
    msg_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    msg_tf.word_wrap = True

    add_footer(slide)
    return slide


def create_slide_3_impact(prs):
    """Slide 3: L'Impact Chiffré"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "À la retraite : jusqu'à 1 600€ de perte mensuelle", 3)

    # Comparaison visuelle
    # Box "Aujourd'hui"
    today_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(4.5), Inches(2.5)
    )
    today_box.fill.solid()
    today_box.fill.fore_color.rgb = GREEN
    today_box.line.fill.background()

    today_title = slide.shapes.add_textbox(
        Inches(1.2), Inches(2), Inches(4.1), Inches(0.5)
    )
    tf = today_title.text_frame
    tf.paragraphs[0].text = "AUJOURD'HUI"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    today_amount = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.8), Inches(4.1), Inches(1)
    )
    tf = today_amount.text_frame
    tf.paragraphs[0].text = "3 000 €"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    today_label = slide.shapes.add_textbox(
        Inches(1.2), Inches(3.7), Inches(4.1), Inches(0.4)
    )
    tf = today_label.text_frame
    tf.paragraphs[0].text = "par mois"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Flèche
    arrow_box = slide.shapes.add_textbox(
        Inches(5.7), Inches(2.6), Inches(1.8), Inches(1)
    )
    tf = arrow_box.text_frame
    tf.paragraphs[0].text = "➜"
    tf.paragraphs[0].font.size = Pt(60)
    tf.paragraphs[0].font.color.rgb = RED
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    arrow_label = slide.shapes.add_textbox(
        Inches(5.5), Inches(3.5), Inches(2.2), Inches(0.5)
    )
    tf = arrow_label.text_frame
    tf.paragraphs[0].text = "-53%"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RED
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Box "Retraite"
    retire_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.5), Inches(2.5)
    )
    retire_box.fill.solid()
    retire_box.fill.fore_color.rgb = RED
    retire_box.line.fill.background()

    retire_title = slide.shapes.add_textbox(
        Inches(8), Inches(2), Inches(4.1), Inches(0.5)
    )
    tf = retire_title.text_frame
    tf.paragraphs[0].text = "À LA RETRAITE"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    retire_amount = slide.shapes.add_textbox(
        Inches(8), Inches(2.8), Inches(4.1), Inches(1)
    )
    tf = retire_amount.text_frame
    tf.paragraphs[0].text = "1 410 €"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    retire_label = slide.shapes.add_textbox(
        Inches(8), Inches(3.7), Inches(4.1), Inches(0.4)
    )
    tf = retire_label.text_frame
    tf.paragraphs[0].text = "par mois"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Calcul de perte
    loss_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(4.6), Inches(7.33), Inches(1.3)
    )
    loss_box.fill.solid()
    loss_box.fill.fore_color.rgb = AUBERGINE
    loss_box.line.fill.background()

    loss_text1 = slide.shapes.add_textbox(
        Inches(3.2), Inches(4.8), Inches(6.93), Inches(0.5)
    )
    tf = loss_text1.text_frame
    tf.paragraphs[0].text = "Perte mensuelle : -1 590 €  |  Perte annuelle : -19 080 €"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Message marquant
    msg_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(6.1), Inches(11.33), Inches(0.9)
    )
    msg_box.fill.solid()
    msg_box.fill.fore_color.rgb = GOLDEN_BROWN
    msg_box.line.fill.background()

    msg_text = slide.shapes.add_textbox(
        Inches(1.2), Inches(6.25), Inches(10.93), Inches(0.6)
    )
    tf = msg_text.text_frame
    tf.paragraphs[
        0
    ].text = (
        "C'est comme perdre un 13ème, 14ème, 15ème... mois de salaire chaque année."
    )
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide


def create_slide_4_solution(prs):
    """Slide 4: La Solution - Le PER"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide, "Le Plan Épargne Retraite : compensez jusqu'à 30% de vos revenus", 4
    )

    # 3 cartes d'avantages
    cards = [
        (
            "💰",
            "Complément de revenus",
            [
                "Reconstituez jusqu'à 30%",
                "de vos revenus actuels",
                "Limitez la perte à 20%",
            ],
        ),
        (
            "📉",
            "Avantages fiscaux",
            ["Versements déductibles", "de vos impôts", "Économie immédiate"],
        ),
        (
            "🔄",
            "Flexibilité totale",
            ["Versements libres", "ou programmés", "Capital, rente ou les deux"],
        ),
    ]

    x_positions = [0.5, 4.6, 8.7]

    for i, (icon, title, bullets) in enumerate(cards):
        x = x_positions[i]

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.8), Inches(3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = GOLDEN_BROWN

        # Icon
        icon_box = slide.shapes.add_textbox(
            Inches(x), Inches(1.6), Inches(3.8), Inches(0.8)
        )
        tf = icon_box.text_frame
        tf.paragraphs[0].text = icon
        tf.paragraphs[0].font.size = Pt(40)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(2.4), Inches(3.6), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = AUBERGINE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Bullets
        y = 3
        for bullet in bullets:
            bullet_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y), Inches(3.4), Inches(0.4)
            )
            tf = bullet_box.text_frame
            tf.paragraphs[0].text = f"• {bullet}"
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = DARK_GRAY
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            y += 0.35

    # Comparaison graphique en bas
    compare_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.8), Inches(12.33), Inches(0.5)
    )
    tf = compare_title.text_frame
    tf.paragraphs[0].text = "Votre situation à la retraite"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = AUBERGINE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Sans PER
    sans_label = slide.shapes.add_textbox(
        Inches(1), Inches(5.4), Inches(2.5), Inches(0.4)
    )
    tf = sans_label.text_frame
    tf.paragraphs[0].text = "Sans PER"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = DARK_GRAY

    sans_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(5.3), Inches(2.2), Inches(0.5)
    )
    sans_bar.fill.solid()
    sans_bar.fill.fore_color.rgb = RED
    sans_bar.line.fill.background()

    sans_pct = slide.shapes.add_textbox(
        Inches(5.8), Inches(5.4), Inches(1), Inches(0.4)
    )
    tf = sans_pct.text_frame
    tf.paragraphs[0].text = "-53%"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RED

    # Avec PER
    avec_label = slide.shapes.add_textbox(
        Inches(1), Inches(6), Inches(2.5), Inches(0.4)
    )
    tf = avec_label.text_frame
    tf.paragraphs[0].text = "Avec PER"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = DARK_GRAY

    avec_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(5.9), Inches(3.8), Inches(0.5)
    )
    avec_bar.fill.solid()
    avec_bar.fill.fore_color.rgb = GREEN
    avec_bar.line.fill.background()

    avec_pct = slide.shapes.add_textbox(Inches(7.4), Inches(6), Inches(1), Inches(0.4))
    tf = avec_pct.text_frame
    tf.paragraphs[0].text = "-20%"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GREEN

    # Économie
    eco_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(5.4), Inches(4), Inches(1)
    )
    eco_box.fill.solid()
    eco_box.fill.fore_color.rgb = GOLDEN_BROWN
    eco_box.line.fill.background()

    eco_text = slide.shapes.add_textbox(
        Inches(8.7), Inches(5.6), Inches(3.6), Inches(0.6)
    )
    tf = eco_text.text_frame
    tf.paragraphs[0].text = "Gain : +33% de revenus conservés"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide


def create_slide_5_offer(prs):
    """Slide 5: L'Offre COGOHR Exclusive"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Adhérents COGOHR : des avantages exclusifs", 5)

    # Badge partenariat
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.2), Inches(4.33), Inches(0.6)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLDEN_BROWN
    badge.line.fill.background()

    badge_text = slide.shapes.add_textbox(
        Inches(4.7), Inches(1.3), Inches(3.93), Inches(0.4)
    )
    tf = badge_text.text_frame
    tf.paragraphs[0].text = "🤝 PARTENAIRE OFFICIEL"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 3 grandes cartes
    advantages = [
        (
            "🎁",
            "GRATUIT",
            "Étude personnalisée",
            [
                "Analyse complète",
                "de votre situation",
                "Simulation personnalisée",
                "Sans aucun engagement",
            ],
        ),
        (
            "0%",
            "0% FRAIS",
            "Frais d'entrée offerts",
            [
                "Économie jusqu'à 3%",
                "du capital investi",
                "Plusieurs milliers d'€",
                "Exclusif COGOHR",
            ],
        ),
        (
            "👥",
            "EXPERT",
            "Accompagnement dédié",
            [
                "Spécialistes fonction",
                "publique hospitalière",
                "Suivi personnalisé",
                "Présence La Réunion",
            ],
        ),
    ]

    x_positions = [0.5, 4.6, 8.7]

    for i, (icon, highlight, title, bullets) in enumerate(advantages):
        x = x_positions[i]

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2), Inches(3.8), Inches(4.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = GOLDEN_BROWN
        card.line.width = Pt(2)

        # Icon circle
        if icon in ["🎁", "👥"]:
            icon_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.2), Inches(3.8), Inches(0.9)
            )
            tf = icon_box.text_frame
            tf.paragraphs[0].text = icon
            tf.paragraphs[0].font.size = Pt(50)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        else:
            # Special styling for 0%
            zero_box = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x + 1.3), Inches(2.2), Inches(1.2), Inches(0.9)
            )
            zero_box.fill.solid()
            zero_box.fill.fore_color.rgb = GREEN
            zero_box.line.fill.background()

            zero_text = slide.shapes.add_textbox(
                Inches(x + 1.3), Inches(2.35), Inches(1.2), Inches(0.6)
            )
            tf = zero_text.text_frame
            tf.paragraphs[0].text = "0%"
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Highlight
        hl_box = slide.shapes.add_textbox(
            Inches(x), Inches(3.2), Inches(3.8), Inches(0.5)
        )
        tf = hl_box.text_frame
        tf.paragraphs[0].text = highlight
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = GREEN if "0%" in highlight else GOLDEN_BROWN
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(3.7), Inches(3.6), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = AUBERGINE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Bullets
        y = 4.3
        for bullet in bullets:
            bullet_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y), Inches(3.4), Inches(0.35)
            )
            tf = bullet_box.text_frame
            tf.paragraphs[0].text = f"✓ {bullet}"
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = DARK_GRAY
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            y += 0.35

    # Message en bas
    msg_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6.7), Inches(9.33), Inches(0.5)
    )
    msg_box.fill.solid()
    msg_box.fill.fore_color.rgb = AUBERGINE
    msg_box.line.fill.background()

    msg_text = slide.shapes.add_textbox(
        Inches(2.2), Inches(6.8), Inches(8.93), Inches(0.3)
    )
    tf = msg_text.text_frame
    tf.paragraphs[
        0
    ].text = "Ces avantages sont réservés exclusivement aux adhérents COGOHR"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide


def create_slide_6_process(prs):
    """Slide 6: Notre Processus en 3 étapes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "3 étapes simples pour sécuriser votre retraite", 6)

    # Timeline
    steps = [
        (
            "1",
            "🔍",
            "DIAGNOSTIC",
            [
                "Analyse de votre",
                "situation actuelle",
                "Simulation retraite",
                "Capacité d'épargne",
            ],
        ),
        (
            "2",
            "💡",
            "SOLUTION",
            [
                "Plan d'épargne",
                "sur-mesure",
                "Optimisation fiscale",
                "Supports adaptés",
            ],
        ),
        (
            "3",
            "📈",
            "SUIVI",
            ["Points réguliers", "Ajustements", "Reporting annuel", "Accompagnement"],
        ),
    ]

    x_positions = [0.7, 4.7, 8.7]

    # Ligne de connexion
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(2.8), Inches(9), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLDEN_BROWN
    line.line.fill.background()

    for i, (num, icon, title, bullets) in enumerate(steps):
        x = x_positions[i]

        # Cercle numéroté
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 1.5), Inches(2.4), Inches(0.9), Inches(0.9)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = GOLDEN_BROWN
        circle.line.fill.background()

        num_text = slide.shapes.add_textbox(
            Inches(x + 1.5), Inches(2.55), Inches(0.9), Inches(0.6)
        )
        tf = num_text.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(3.5),
            Inches(3.9),
            Inches(3.2),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = GOLDEN_BROWN

        # Icon
        icon_box = slide.shapes.add_textbox(
            Inches(x), Inches(3.6), Inches(3.9), Inches(0.7)
        )
        tf = icon_box.text_frame
        tf.paragraphs[0].text = icon
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(4.3), Inches(3.7), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = AUBERGINE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Bullets
        y = 4.9
        for bullet in bullets:
            bullet_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y), Inches(3.5), Inches(0.4)
            )
            tf = bullet_box.text_frame
            tf.paragraphs[0].text = f"• {bullet}"
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = DARK_GRAY
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            y += 0.4

    add_footer(slide)
    return slide


def create_slide_7_cta(prs):
    """Slide 7: Conclusion et CTA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond coloré
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = AUBERGINE
    bg.line.fill.background()

    # Bande dorée en haut
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.15)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = GOLDEN_BROWN
    band.line.fill.background()

    # Titre
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1), Inches(12.33), Inches(1)
    )
    tf = title_box.text_frame
    tf.paragraphs[0].text = "Agissez maintenant"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Sous-titre
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.9), Inches(12.33), Inches(0.6)
    )
    tf = sub_box.text_frame
    tf.paragraphs[0].text = "C'est gratuit et sans engagement"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = GOLDEN_BROWN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Box contact
    contact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(2.8), Inches(7.33), Inches(2.5)
    )
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = WHITE
    contact_box.line.fill.background()

    # Téléphone
    phone_box = slide.shapes.add_textbox(
        Inches(3.2), Inches(3), Inches(6.93), Inches(0.6)
    )
    tf = phone_box.text_frame
    tf.paragraphs[0].text = "📞 05 33 89 14 00"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = AUBERGINE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Email
    email_box = slide.shapes.add_textbox(
        Inches(3.2), Inches(3.7), Inches(6.93), Inches(0.5)
    )
    tf = email_box.text_frame
    tf.paragraphs[0].text = "✉️ contact-reunion@argamconseils.com"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = AUBERGINE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Garantie
    garantie_box = slide.shapes.add_textbox(
        Inches(3.2), Inches(4.5), Inches(6.93), Inches(0.5)
    )
    tf = garantie_box.text_frame
    tf.paragraphs[0].text = "⏱️ Réponse garantie sous 48h"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = GOLDEN_BROWN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # CTA
    cta_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.5), Inches(8.33), Inches(0.9)
    )
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = GOLDEN_BROWN
    cta_box.line.fill.background()

    cta_text = slide.shapes.add_textbox(
        Inches(2.7), Inches(5.7), Inches(7.93), Inches(0.5)
    )
    tf = cta_text.text_frame
    tf.paragraphs[0].text = "Demandez votre étude personnalisée gratuite"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Badges partenaires
    badge_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.5)
    )
    tf = badge_box.text_frame
    tf.paragraphs[0].text = "Argam Conseils x COGOHR | Partenaire Officiel"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def main():
    """Fonction principale de création de la présentation"""
    # Créer la présentation
    prs = Presentation()

    # Définir le format 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Création de la présentation COGOHR...")

    # Slide 1: Introduction
    print("  → Slide 1: Introduction")
    add_title_slide(
        prs, "Argam Conseils x COGOHR", "Protégez votre niveau de vie à la retraite"
    )

    # Slide 2: Le Problème
    print("  → Slide 2: Le Problème")
    create_slide_2_problem(prs)

    # Slide 3: L'Impact Chiffré
    print("  → Slide 3: L'Impact Chiffré")
    create_slide_3_impact(prs)

    # Slide 4: La Solution - Le PER
    print("  → Slide 4: La Solution")
    create_slide_4_solution(prs)

    # Slide 5: L'Offre COGOHR Exclusive
    print("  → Slide 5: L'Offre COGOHR")
    create_slide_5_offer(prs)

    # Slide 6: Notre Processus
    print("  → Slide 6: Notre Processus")
    create_slide_6_process(prs)

    # Slide 7: Conclusion et CTA
    print("  → Slide 7: Conclusion")
    create_slide_7_cta(prs)

    # Sauvegarder
    output_path = (
        "/Users/julianbonne/Desktop/ARGAMCONSEIL NEW/Presentation_COGOHR_Argam.pptx"
    )
    prs.save(output_path)
    print(f"\n✅ Présentation créée avec succès !")
    print(f"   Fichier : {output_path}")
    print(f"   Nombre de slides : 7")


if __name__ == "__main__":
    main()
